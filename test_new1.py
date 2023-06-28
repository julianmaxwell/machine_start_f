__author__ = "boss"
__date__ = '2022/7/18 16:03'

import pptx
import PIL
from pptx.util import Inches
from pptx import Presentation
import datetime
import time
import pandas
pptx_file = Presentation()
slide = pptx_file.slides.add_slide(pptx_file.slide_layouts[0])
pptx_file.save(r"C:\Users\Administrator\Desktop\2022年工作文件夹\企业资质202206\公司资质\企业资质202208.pptx")
save_path = r"C:\Users\Administrator\Desktop\2022年工作文件夹\企业资质202206\公司资质"
certificate_use_recorder_file = "certificate_use_recorder_.csv"
import os
from PIL import Image, ImageDraw, ImageFont

def deputy_safe_file(save_path=None, text_=None, user="张山", filter_mark="测试使用",
              source_file=r"C:\Users\Administrator\Desktop\2022年工作文件夹\企业资质202206\公司资质\图片草稿\安许副本_2025_8_10.jpg"):

    print(save_path)
    print(source_file)
    read_image = Image.open(source_file)
    read_image.thumbnail((6600, 4600))
    read_image = read_image.rotate(12, fillcolor=None)
    draw_image = ImageDraw.Draw(read_image)
    fun_t = ImageFont.truetype('simhei',200)
    if text_:
        pass
    else:
        text_ = "企业资质使用 说明"
    draw_image.text((900, 2200), text=text_, font=fun_t)
    draw_image.text((900, 3200), text=text_, font=fun_t)
    draw_image.text((3500, 2700), text=text_, font=fun_t)
    draw_image.text((3500, 1700), text=text_, font=fun_t)
    read_image = read_image.rotate(-12, fillcolor=None)
    read_image.save(r"C:\Users\Administrator\Desktop\2022年工作文件夹\企业资质202206\公司资质\图片草稿\安许副本2_2025_8_10.jpg")
    read_image = read_image.crop(box=(417,569, 5397, 4197))
    use_date = str(pandas.to_datetime(datetime.datetime.now()))
    file_name = "安许副本.jpg"
    path_ = os.path.join(save_path, filter_mark)
    if not os.path.exists(path_):
        os.makedirs(path_)
    save_path_file = os.path.join(path_, file_name)
    print(save_path_file)
    if save_path:
        read_image.save(save_path_file)
        data_record = pandas.DataFrame([[user, filter_mark, use_date, text_]],
                                       columns=["使用人", "查询关键字", "使用时间", "注释文本"])
        if os.path.exists(certificate_use_recorder_file):
            data_record.to_csv(header=False, encoding="gbk", index=False)
        else:
            data_record.to_csv(header=True, encoding="gbk", index=False)
    read_image.show()

# safe_file(save_path, text_="遂宁中环线过街非开挖使用 ", user="boss", filter_mark="遂宁中环线过街非开挖使用")
def main_safe_file(save_path=None, text_=None, user="张山", filter_mark="测试公司使用",
              source_file=r"C:\Users\Administrator\Desktop\2022年工作文件夹\企业资质202206\公司资质\图片草稿\安许正本_2025_8_10.jpg"):
    read_image = Image.open(source_file)
    # read_image.thumbnail((6600, 4600))
    angle = 10
    read_image = read_image.rotate(angle, fillcolor=None)

    draw_image = ImageDraw.Draw(read_image)
    fun_t = ImageFont.truetype('simhei',200)
    if text_:
        pass
    else:
        text_ = "企业资质使用 说明"

    deltax = 200
    delta_y = 500
    draw_image.text((900+deltax, 2600+delta_y), text=text_, font=fun_t)
    draw_image.text((900+deltax, 3600+delta_y), text=text_, font=fun_t)
    draw_image.text((3500+deltax, 3100+delta_y), text=text_, font=fun_t)
    draw_image.text((3500+deltax, 2100+delta_y), text=text_, font=fun_t)
    read_image = read_image.rotate(-angle, fillcolor=None)
    print(read_image.size)
    read_image.save(r"C:\Users\Administrator\Desktop\2022年工作文件夹\企业资质202206\公司资质\图片草稿\安许副本2_2025_8_10.jpg")
    read_image = read_image.crop(box=(350, 800, 7600, 5561))
    use_date = str(pandas.to_datetime(datetime.datetime.now()))
    file_name = "安许正本.jpg"
    path_ = os.path.join(save_path, filter_mark)
    if not os.path.exists(path_):
        os.makedirs(path_)
    save_path_file = os.path.join(path_, file_name)
    if save_path:
        read_image.save(save_path_file)
        data_record = pandas.DataFrame([[user, filter_mark, use_date, text_]],
                                       columns=["使用人", "查询关键字", "使用时间", "注释文本"])
        if os.path.exists(certificate_use_recorder_file):
            data_record.to_csv(header=False, encoding="gbk", index=False)
        else:
            data_record.to_csv(header=True, encoding="gbk", index=False)
    read_image.save(r"C:\Users\Administrator\Desktop\2022年工作文件夹\企业资质202206\公司资质\图片草稿\安许正本2_2025_8_10.jpg")
    read_image.show()

def deputy_municipal_file(save_path=None, text_=None, user="张山", filter_mark="测试公司使用",
              source_file=r"C:\Users\Administrator\Desktop\2022年工作文件夹\企业资质202206\公司资质\图片草稿\市政三级副本2025_4_30.jpg"):
    read_image = Image.open(source_file)
    # read_image.thumbnail((6600, 4600))
    angle = 5
    read_image = read_image.rotate(angle, fillcolor=None)

    draw_image = ImageDraw.Draw(read_image)
    fun_t = ImageFont.truetype('simhei',200)
    if text_:
        pass
    else:
        text_ = "企业资质使用 说明"

    deltax = 200
    delta_y = 350
    draw_image.text((900+deltax, 2600+delta_y), text=text_, font=fun_t)
    draw_image.text((900+deltax, 3600+delta_y), text=text_, font=fun_t)
    draw_image.text((3500+deltax, 3100+delta_y), text=text_, font=fun_t)
    draw_image.text((3500+deltax, 2100+delta_y), text=text_, font=fun_t)
    read_image = read_image.rotate(-angle, fillcolor=None)
    read_image = read_image.crop(box=(377, 313, 5449, 7513))
    use_date = str(pandas.to_datetime(datetime.datetime.now()))
    file_name = "市政施工副本.jpg"
    path_ = os.path.join(save_path, filter_mark)
    if not os.path.exists(path_):
        os.makedirs(path_)
    save_path_file = os.path.join(path_, file_name)
    if save_path:
        read_image.save(save_path_file)
        data_record = pandas.DataFrame([[user, filter_mark, use_date, text_]],
                                       columns=["使用人", "查询关键字", "使用时间", "注释文本"])
        if os.path.exists(certificate_use_recorder_file):
            data_record.to_csv(header=False, encoding="gbk", index=False)
        else:
            data_record.to_csv(header=True, encoding="gbk", index=False)
    read_image.save(r"C:\Users\Administrator\Desktop\2022年工作文件夹\企业资质202206\公司资质\图片草稿\市政三级2副本2025_4_30.jpg")

    read_image.show()

def main_municipal_file(save_path=None, text_=None, user="张山", filter_mark="测试公司使用",
              source_file=r"C:\Users\Administrator\Desktop\2022年工作文件夹\企业资质202206\公司资质\图片草稿\市政三级正本2025_4_30.jpg"):
    read_image = Image.open(source_file)
    # read_image.thumbnail((6600, 4600))
    angle = 5
    read_image = read_image.rotate(angle, fillcolor=None)

    draw_image = ImageDraw.Draw(read_image)
    fun_t = ImageFont.truetype('simhei',200)
    if text_:
        pass
    else:
        text_ = "企业资质使用 说明"

    deltax = 200
    delta_y = 350
    draw_image.text((900+deltax, 2600+delta_y), text=text_, font=fun_t)
    draw_image.text((900+deltax, 3600+delta_y), text=text_, font=fun_t)
    draw_image.text((3500+deltax, 3100+delta_y), text=text_, font=fun_t)
    draw_image.text((3500+deltax, 2100+delta_y), text=text_, font=fun_t)
    read_image = read_image.rotate(-angle, fillcolor=None)
    read_image = read_image.crop(box=(377, 313, 5449, 7513))
    use_date = str(pandas.to_datetime(datetime.datetime.now()))
    file_name = "市政施工正本.jpg"
    path_ = os.path.join(save_path, filter_mark)
    if not os.path.exists(path_):
        os.makedirs(path_)
    save_path_file = os.path.join(path_, file_name)
    if save_path:
        read_image.save(save_path_file)
        data_record = pandas.DataFrame([[user, filter_mark, use_date, text_]],
                                       columns=["使用人", "查询关键字", "使用时间", "注释文本"])
        if os.path.exists(certificate_use_recorder_file):
            data_record.to_csv(header=False, encoding="gbk", index=False)
        else:
            data_record.to_csv(header=True, encoding="gbk", index=False)
    # read_image.save(r"C:\Users\Administrator\Desktop\2022年工作文件夹\企业资质202206\公司资质\图片草稿\市政三级2副本2025_4_30.jpg")

    read_image.show()


def deputy_business(save_path=None, text_=None, user="张山", filter_mark="测试公司使用",
                            source_file=r"C:\Users\Administrator\Desktop\2022年工作文件夹\企业资质202206\公司资质\图片草稿\营业执照副本.jpg"):
        read_image = Image.open(source_file)
        # read_image.thumbnail((6600, 4600))
        angle = 5
        read_image = read_image.rotate(angle, fillcolor=None)

        draw_image = ImageDraw.Draw(read_image)
        fun_t = ImageFont.truetype('simhei', 200)
        if text_:
            pass
        else:
            text_ = "企业资质使用 说明"

        deltax = 200
        delta_y = 1230
        draw_image.text((900 + deltax, 2600 + delta_y), text=text_, font=fun_t)
        draw_image.text((900 + deltax, 3600 + delta_y), text=text_, font=fun_t)
        draw_image.text((3500 + deltax, 3100 + delta_y), text=text_, font=fun_t)
        draw_image.text((3500 + deltax, 2100 + delta_y), text=text_, font=fun_t)
        read_image = read_image.rotate(-angle, fillcolor=None)
        print(read_image.size)
        read_image = read_image.crop(box=(400, 650, 5449, 7600))
        use_date = str(pandas.to_datetime(datetime.datetime.now()))
        file_name = "营业执照副本.jpg"
        path_ = os.path.join(save_path, filter_mark)
        if not os.path.exists(path_):
            os.makedirs(path_)
        save_path_file = os.path.join(path_, file_name)
        if save_path:
            read_image.save(save_path_file)
            data_record = pandas.DataFrame([[user, filter_mark, use_date, text_]],
                                           columns=["使用人", "查询关键字", "使用时间", "注释文本"])
            if os.path.exists(certificate_use_recorder_file):
                data_record.to_csv(header=False, encoding="gbk", index=False)
            else:
                data_record.to_csv(header=True, encoding="gbk", index=False)
        # read_image.save(r"C:\Users\Administrator\Desktop\2022年工作文件夹\企业资质202206\公司资质\图片草稿\市政三级2副本2025_4_30.jpg")

        read_image.show()


def deputy_communication_certificate(save_path=None, text_=None, user="张山", filter_mark="测试公司使用",
                    source_file=r"C:\Users\Administrator\Desktop\2022年工作文件夹\企业资质202206\公司资质\图片草稿\通信副本2022_12_31.jpg"):
    read_image = Image.open(source_file)
    # read_image.thumbnail((6600, 4600))
    angle = 7
    read_image = read_image.rotate(angle, fillcolor=None)

    draw_image = ImageDraw.Draw(read_image)
    fun_t = ImageFont.truetype('simhei', 200)
    if text_:
        pass
    else:
        text_ = "企业资质使用 说明"

    deltax = 200
    delta_y = 550
    draw_image.text((900 + deltax, 2600 + delta_y), text=text_, font=fun_t)
    draw_image.text((900 + deltax, 3600 + delta_y), text=text_, font=fun_t)
    draw_image.text((3500 + deltax, 3100 + delta_y), text=text_, font=fun_t)
    draw_image.text((3500 + deltax, 2100 + delta_y), text=text_, font=fun_t)
    read_image = read_image.rotate(-angle, fillcolor=None)
    print(read_image.size)
    read_image = read_image.crop(box=(590, 650, 5200, 7100))
    use_date = str(pandas.to_datetime(datetime.datetime.now()))
    file_name = "通信施工副本.jpg"
    path_ = os.path.join(save_path, filter_mark)
    if not os.path.exists(path_):
        os.makedirs(path_)
    save_path_file = os.path.join(path_, file_name)
    if save_path:
        read_image.save(save_path_file)
        data_record = pandas.DataFrame([[user, filter_mark, use_date, text_]],
                                       columns=["使用人", "查询关键字", "使用时间", "注释文本"])
        if os.path.exists(certificate_use_recorder_file):
            data_record.to_csv(header=False, encoding="gbk", index=False)
        else:
            data_record.to_csv(header=True, encoding="gbk", index=False)
    read_image.save(r"C:\Users\Administrator\Desktop\2022年工作文件夹\企业资质202206\公司资质\图片草稿\市政三级2副本2025_4_30.jpg")

    read_image.show()



def make_certificate_(save_path=None,
                      text_=None,
                      user="张山",
                      filter_mark="测试公司使用",
                    source_file_path=r"C:\Users\Administrator\Desktop\2022年工作文件夹\企业资质202206\公司资质\图片草稿"):
    file_name = {
      "deputy_safe_file":  "安许副本_2025_8_10.jpg",
      "main_safe_file":  "安许正本_2025_8_10.jpg",
      "deputy_municipal_file":  "市政三级副本2025_4_30.jpg",
      "main_municipal_file":  "市政三级正本2025_4_30.jpg",
      "deputy_business":  "营业执照副本.jpg",
      "deputy_communication_certificate":  "通信副本2022_12_31.jpg",
    }
    save_pa = {}
    for xx in file_name.keys():
        file_name[xx] = os.path.join(source_file_path, file_name[xx])


    deputy_safe_file(save_path, text_=text_, user=user, filter_mark=filter_mark, source_file=file_name["deputy_safe_file"])
    main_safe_file(save_path, text_=text_, user=user, filter_mark=filter_mark, source_file=file_name["main_safe_file"])
    deputy_municipal_file(save_path, text_=text_, user=user, filter_mark=filter_mark, source_file=file_name["deputy_municipal_file"])
    main_municipal_file(save_path, text_=text_, user=user, filter_mark=filter_mark, source_file=file_name["main_municipal_file"])
    deputy_business(save_path, text_=text_, user=user, filter_mark=filter_mark, source_file=file_name["deputy_business"])
    deputy_communication_certificate(save_path, text_=text_, user=user, filter_mark=filter_mark, source_file=file_name["deputy_communication_certificate"])


make_certificate_(save_path=r"C:\Users\Administrator\Desktop\2022年工作文件夹\企业资质使用档案",
                  text_="2022年遂宁中环线过街非开挖",
                  user="boss",
                  filter_mark="2022年遂宁中环线过街非开挖",
                  source_file_path=r"C:\Users\Administrator\Desktop\2022年工作文件夹\企业资质202206\公司资质\图片草稿")